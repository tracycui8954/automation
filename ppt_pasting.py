#install library 
#pip install mss

##Take screenshot of part of the screen
import mss
import mss.tools
from datetime import datetime

with mss.mss() as sct:
    # Get information of monitor 2
    monitor_number = 2
    mon = sct.monitors[monitor_number]
    # The part of screen to capture, depends on the size of your dashboard, try several times to get the perfect parameters
    # Format all dashboards to have the same size for the script to work accorss differnet workbooks
    monitor = {
        "top": mon["top"] + 115,  # 115px from the top
        "left": mon["left"] + 320,  # 260px from the left
        "width": 1300,              #1300px width
        "height": 700,              #700px height
        "mon": monitor_number,
    }
    # Each screensot is saved with a timestamp in its name
    dt = datetime.now()
    fname = "pic_{}.png".format(dt.strftime("%H%M_%S"))
    output = "C:\\Users\\Desktop\\Screenshot\\" + fname
    # Grab the screenshot
    sct_img = sct.grab(monitor)
    # Save to the picture file
    mss.tools.to_png(sct_img.rgb, sct_img.size, output=output)
    
#Install the library
#pip install python-pptx

##Paste images into ppt and save as a ppt###
import pptx
from pptx import Presentation
import pptx.util
import glob
import scipy.misc

#Text to show on the tile slide
OUTPUT_TAG = "Data Review for Facebook"

# paste in a new powerpoint
prs = pptx.Presentation()
# paste in an existing powerpoint
# prs_exists = pptx.Presentation("Existing.pptx")

# Set the width of slides

# default slide width
#prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
#title = slide.shapes.title
#title.text = OUTPUT_TAG

pic_left  = int(prs.slide_width * 0.1)
pic_top   = int(prs.slide_height * 0.2)
pic_width = int(prs.slide_width * 0.8)


for g in glob.glob('C:/users/Desktop/Screenshot/*'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(50, 50, prs.slide_width,50)  #(left, top, width, height)
    p = tb.text_frame.add_paragraph()
    p.text = "Slide title you would like to input for every slide"
    p.font.size = pptx.util.Pt(15)
    p.font.bold = True
    img = scipy.misc.imread(g)
    pic_height = int(pic_width * img.shape[0] / img.shape[1])
    pic = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
    
#Save as the name you want    
prs.save('Facebook_DataReview.pptx')  
